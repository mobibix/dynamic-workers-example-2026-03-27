"""
Extract person names and photos from PowerPoint files.

Each slide is expected to have:
  - One photo (the largest picture shape is used)
  - One text element containing the person's name

Usage:
    python extract_pptx.py staff.pptx
    python extract_pptx.py slides1.pptx slides2.pptx --output-dir data
"""

import argparse
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def sanitize_filename(name: str) -> str:
    safe = "".join(c for c in name if c.isalnum() or c in (" ", "-", "_")).strip()
    return safe.replace(" ", "_")


def extract_slide(slide, slide_num: int):
    """Return (name, image_blob, image_ext) or None if slide can't be parsed."""
    name = None
    best_picture = None
    best_picture_area = 0

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Use first non-trivial text block as the name
            if text and 1 < len(text) < 80 and name is None:
                name = text

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            area = shape.width * shape.height
            if area > best_picture_area:
                best_picture_area = area
                best_picture = shape

    if not name:
        print(f"  Slide {slide_num}: skipped — no name text found")
        return None
    if not best_picture:
        print(f"  Slide {slide_num}: skipped — no picture found")
        return None

    ext = best_picture.image.ext or "png"
    return name, best_picture.image.blob, ext


def extract_people(pptx_paths, output_dir: str = "data"):
    images_dir = Path(output_dir) / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    people = []
    seen_names: set = set()

    for pptx_path in pptx_paths:
        print(f"\nProcessing: {pptx_path}")
        try:
            prs = Presentation(pptx_path)
        except Exception as exc:
            print(f"  ERROR opening file: {exc}")
            continue

        for slide_num, slide in enumerate(prs.slides, start=1):
            result = extract_slide(slide, slide_num)
            if result is None:
                continue

            name, image_blob, image_ext = result

            if name in seen_names:
                print(f"  Slide {slide_num}: skipped duplicate '{name}'")
                continue
            seen_names.add(name)

            safe = sanitize_filename(name)
            filename = f"{safe}.{image_ext}"
            # Avoid collisions from different names that sanitize to the same string
            dest = images_dir / filename
            counter = 1
            while dest.exists():
                filename = f"{safe}_{counter}.{image_ext}"
                dest = images_dir / filename
                counter += 1

            dest.write_bytes(image_blob)

            people.append({"name": name, "image_path": str(Path("images") / filename)})
            print(f"  Slide {slide_num}: extracted '{name}'")

    index_path = Path(output_dir) / "people.json"
    with open(index_path, "w", encoding="utf-8") as f:
        json.dump(people, f, indent=2, ensure_ascii=False)

    print(f"\nDone — {len(people)} people saved to '{output_dir}/'")
    print(f"Run  streamlit run app.py  to start the trainer.")
    return people


def main():
    parser = argparse.ArgumentParser(
        description="Extract faces and names from PowerPoint files"
    )
    parser.add_argument("pptx_files", nargs="+", help="PowerPoint file(s) to process")
    parser.add_argument(
        "--output-dir",
        default="data",
        help="Directory to write images and people.json (default: data)",
    )
    args = parser.parse_args()

    missing = [p for p in args.pptx_files if not os.path.exists(p)]
    if missing:
        for m in missing:
            print(f"File not found: {m}", file=sys.stderr)
        sys.exit(1)

    extract_people(args.pptx_files, args.output_dir)


if __name__ == "__main__":
    main()
